#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception as exc:  # pragma: no cover
    raise SystemExit(
        "python-pptx is required. Install with: pip install python-pptx"
    ) from exc


# NOT: Asagidaki "V2" etiketleri surum fosili degildir; bu betigin amaci
# yatirimci sunumuna (PPTX) bilerek "V2" damgasini ISLEMEKtir, yani bunlar
# betigin fonksiyonel yuku/payload'udur. Surum kimligi tek bir kaynaktan
# (release_manifest/config) turetilmedigi icin, build numarasi degisirse bu
# string'ler elle guncellenmelidir (tek-kaynak baglama yapilmadi).
REPLACEMENTS = [
    ("Build 30", "Build 30 V2"),
    ("Build30", "Build30 V2"),
    ("V30.0", "V30.0-V2"),
]

V2_BULLETS_EN = [
    "Dedup pipeline enabled (global/stage scope)",
    "MoE parallel dispatch path (sequential fallback retained)",
    "Liquid/CfC fast-path opt-in for runtime efficiency",
    "Training gates + SOP evidence hardened",
    "Claim boundary unchanged (pre-training, no benchmark claims)",
]

V2_BULLETS_TR = [
    "Dedup pipeline aktif (global/sahne kapsamı)",
    "MoE paralel dispatch yolu (sequential fallback korunur)",
    "Liquid/CfC fast-path opt-in (verimlilik)",
    "Eğitim kapıları + SOP kanıtları güçlendirildi",
    "İddia sınırı aynı (eğitim öncesi, benchmark yok)",
]


def _replace_text(text: str) -> str:
    out = text
    for old, new in REPLACEMENTS:
        out = out.replace(old, new)
    return out


def _iter_text_shapes(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            yield slide, shape


def _find_title_slide(prs: Presentation, title: str):
    for slide in prs.slides:
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            if title_shape.text.strip() == title:
                return slide
    return None


def _safe_layout(prs: Presentation):
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    if len(prs.slide_layouts) > 0:
        return prs.slide_layouts[0]
    return None


def _ensure_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def _update_or_add_v2_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = _find_title_slide(prs, title)
    if slide is None:
        layout = _safe_layout(prs)
        if layout is None:
            raise RuntimeError("No slide layouts available in PPTX")
        slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        slide.shapes.title.text = title
    else:
        tf = _ensure_textbox(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
        tf.text = title

    # Body
    body = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body = shape.text_frame
            break
    if body is None:
        body = _ensure_textbox(slide, Inches(0.8), Inches(1.6), Inches(12.0), Inches(4.5))
    body.clear()
    body.text = bullets[0]
    for line in bullets[1:]:
        p = body.add_paragraph()
        p.text = line
        p.level = 0


def update_pptx(path: Path, language: str, out: Path | None = None) -> Path:
    prs = Presentation(str(path))

    for _, shape in _iter_text_shapes(prs):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.text = _replace_text(run.text)

    if language == "tr":
        _update_or_add_v2_slide(prs, "Build 30 V2 Güncellemesi", V2_BULLETS_TR)
    else:
        _update_or_add_v2_slide(prs, "Build 30 V2 Update", V2_BULLETS_EN)

    target = out or path
    prs.save(str(target))
    return target


def main() -> int:
    parser = argparse.ArgumentParser(description="Update investor deck PPTX to Build 30 V2.")
    parser.add_argument("--input", default="reports/investor_deck.pptx")
    parser.add_argument("--input-tr", default="reports/investor_deck_TR.pptx")
    parser.add_argument("--out", default="")
    parser.add_argument("--out-tr", default="")
    args = parser.parse_args()

    inputs = [
        (Path(args.input), "en", Path(args.out) if args.out else None),
        (Path(args.input_tr), "tr", Path(args.out_tr) if args.out_tr else None),
    ]

    for path, lang, out in inputs:
        if not path.exists():
            print(f"[skip] missing: {path}")
            continue
        target = update_pptx(path, lang, out=out)
        print(f"[ok] updated: {target}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
